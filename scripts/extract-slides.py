#!/usr/bin/env python3
"""
Extract text content from lecture slides (PDF, PPTX, DOCX).
Usage: python extract-slides.py <file_path>
Output: extracted text to stdout (markdown-ish format)
"""

import sys
import os
import io

# Fix Windows encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')


def extract_pdf(path):
    """Extract text from PDF using PyPDF2."""
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    lines = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            lines.append(f"--- Slide/Page {i+1} ---")
            lines.append(text.strip())
            lines.append("")
    return "\n".join(lines)


def extract_pptx(path):
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Preserve bullet level
                        indent = "  " * para.level if para.level else ""
                        slide_texts.append(f"{indent}{text}")
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append("| " + " | ".join(cells) + " |")
        if slide_texts:
            lines.append(f"--- Slide {i+1} ---")
            lines.append("\n".join(slide_texts))
            lines.append("")
    return "\n".join(lines)


def extract_docx(path):
    """Extract text from DOCX using python-docx."""
    from docx import Document
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract-slides.py <file_path>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.exists(path):
        print(f"Error: file not found: {path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()

    try:
        if ext == ".pdf":
            print(extract_pdf(path))
        elif ext in (".pptx", ".ppt"):
            print(extract_pptx(path))
        elif ext in (".docx", ".doc"):
            print(extract_docx(path))
        else:
            print(f"Error: unsupported format: {ext}", file=sys.stderr)
            sys.exit(1)
    except Exception as e:
        print(f"Error extracting {path}: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
